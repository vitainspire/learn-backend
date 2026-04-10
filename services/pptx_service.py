import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime


# ---------------------------------------------------------------------------
# Design constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# White background, dark text — printable
BG          = RGBColor(255, 255, 255)
TEXT_DARK   = RGBColor(30,  30,  30)
TEXT_MUTED  = RGBColor(100, 100, 100)
TEXT_WHITE  = RGBColor(255, 255, 255)

# One accent colour per 5E phase
PHASE_META = {
    "cover":     {"color": RGBColor(26,  35,  126), "number": "",   "label": "",           "sub": ""},
    "engage":    {"color": RGBColor(230, 81,  0),   "number": "01", "label": "ENGAGE",     "sub": "The Hook"},
    "explore":   {"color": RGBColor(27,  94,  32),  "number": "02", "label": "EXPLORE",    "sub": "Student Discovery"},
    "explain":   {"color": RGBColor(13,  71,  161), "number": "03", "label": "EXPLAIN",    "sub": "Core Concepts"},
    "elaborate": {"color": RGBColor(74,  20,  140), "number": "04", "label": "ELABORATE",  "sub": "Guided & Independent Practice"},
    "evaluate":  {"color": RGBColor(183, 28,  28),  "number": "05", "label": "EVALUATE",   "sub": "Check for Understanding"},
    "closure":   {"color": RGBColor(0,   77,  64),  "number": "06", "label": "CLOSURE",    "sub": "Wrap-Up & Bridge"},
}

HEADER_H   = Inches(2.2)   # coloured band height
MARGIN     = Inches(0.55)
BODY_TOP   = HEADER_H + Inches(0.35)
BODY_H     = SLIDE_H - BODY_TOP - Inches(0.3)
BODY_W     = SLIDE_W - MARGIN * 2


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs):
    layout = prs.slide_layouts[6]   # fully blank
    slide  = prs.slides.add_slide(layout)
    bg     = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def _header_band(slide, phase_key):
    """Coloured top band with phase number + name + subtitle."""
    meta  = PHASE_META[phase_key]
    color = meta["color"]

    # Full-width colour rectangle
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, HEADER_H)
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()

    # Phase number (top-left, small)
    if meta["number"]:
        nb = slide.shapes.add_textbox(MARGIN, Inches(0.25), Inches(1), Inches(0.5))
        p  = nb.text_frame.paragraphs[0]
        p.text = meta["number"]
        p.font.size  = Pt(14)
        p.font.bold  = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.color.rgb.red   = 255
        p.font.color.rgb.green = 255
        p.font.color.rgb.blue  = 255
        # use a lighter tint via alpha trick — just white with opacity
        nb.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE

    # Phase label (large, bold)
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.55), BODY_W, Inches(1.0))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.text            = meta["label"]
    p.font.bold       = True
    p.font.size       = Pt(44)
    p.font.color.rgb  = TEXT_WHITE

    # Subtitle
    tb2 = slide.shapes.add_textbox(MARGIN, Inches(1.55), BODY_W, Inches(0.5))
    p2  = tb2.text_frame.paragraphs[0]
    p2.text           = meta["sub"]
    p2.font.size      = Pt(16)
    p2.font.color.rgb = RGBColor(220, 220, 220)

    # Thin white rule below band
    rule = slide.shapes.add_shape(6, Inches(0), HEADER_H, SLIDE_W, Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()


def _body_textbox(slide, text_lines: list[tuple], top_offset=Inches(0)):
    """
    text_lines: list of (text, font_size, bold, color) tuples.
    Returns the textbox so the caller can add more paragraphs.
    """
    tb = slide.shapes.add_textbox(
        MARGIN,
        BODY_TOP + top_offset,
        BODY_W,
        BODY_H - top_offset,
    )
    tf        = tb.text_frame
    tf.word_wrap = True

    first = True
    for text, size, bold, color in text_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first       = False
        p.text      = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after    = Pt(8)

    return tb


def _bullet_lines(items: list[str], prefix="•") -> list[tuple]:
    """Convert a list of strings into body_textbox tuples."""
    return [(f"{prefix}  {item}", 18, False, TEXT_DARK) for item in items if item]


def _truncate(text: str, limit=280) -> str:
    if not text:
        return ""
    return (text[:limit] + "…") if len(text) > limit else text


# ---------------------------------------------------------------------------
# One function per flashcard
# ---------------------------------------------------------------------------

def _slide_cover(prs, lesson_data: dict):
    slide = _blank_slide(prs)
    meta  = lesson_data.get("meta", {})

    # Top colour band (cover uses deep navy)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(3.2))
    band.fill.solid()
    band.fill.fore_color.rgb = PHASE_META["cover"]["color"]
    band.line.fill.background()

    # "LESSON FLASHCARDS" label
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.4), BODY_W, Inches(0.6))
    p  = tb.text_frame.paragraphs[0]
    p.text           = "LESSON FLASHCARDS  ·  5E MODEL"
    p.font.size      = Pt(13)
    p.font.bold      = True
    p.font.color.rgb = RGBColor(180, 180, 220)

    # Lesson title
    tb2 = slide.shapes.add_textbox(MARGIN, Inches(1.0), BODY_W, Inches(1.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2  = tf2.paragraphs[0]
    p2.text           = meta.get("lesson_title", "Untitled Lesson")
    p2.font.bold      = True
    p2.font.size      = Pt(40)
    p2.font.color.rgb = TEXT_WHITE

    # Grade / Duration
    tb3 = slide.shapes.add_textbox(MARGIN, Inches(3.5), BODY_W, Inches(0.6))
    p3  = tb3.text_frame.paragraphs[0]
    grade    = meta.get("grade", "")
    duration = meta.get("duration", "")
    p3.text           = f"Grade {grade}   ·   {duration}"
    p3.font.size      = Pt(18)
    p3.font.color.rgb = TEXT_MUTED

    # Objectives
    objectives = lesson_data.get("objective", [])
    if objectives:
        tb4 = slide.shapes.add_textbox(MARGIN, Inches(4.2), BODY_W, Inches(2.8))
        tf4 = tb4.text_frame
        tf4.word_wrap = True
        header        = tf4.paragraphs[0]
        header.text           = "LEARNING OBJECTIVES"
        header.font.size      = Pt(12)
        header.font.bold      = True
        header.font.color.rgb = TEXT_MUTED
        header.space_after    = Pt(6)
        for obj in objectives[:4]:
            p = tf4.add_paragraph()
            p.text           = f"•  {obj}"
            p.font.size      = Pt(15)
            p.font.color.rgb = TEXT_DARK
            p.space_after    = Pt(4)


def _slide_engage(prs, lesson_data: dict):
    engage = lesson_data.get("engage") or {}
    slide  = _blank_slide(prs)
    _header_band(slide, "engage")

    activity = engage.get("activity", "No activity provided.")
    duration = engage.get("duration", "")

    lines = [
        ("ACTIVITY", 12, True, TEXT_MUTED),
        (_truncate(activity, 400), 20, False, TEXT_DARK),
    ]
    if duration:
        lines += [("", 10, False, TEXT_DARK), (f"⏱  {duration}", 14, False, TEXT_MUTED)]
    _body_textbox(slide, lines)


def _slide_explore(prs, lesson_data: dict):
    explore = lesson_data.get("explore") or {}
    slide   = _blank_slide(prs)
    _header_band(slide, "explore")

    activity = explore.get("activity", "No activity provided.")
    duration = explore.get("duration", "")

    lines = [
        ("STUDENT ACTIVITY", 12, True, TEXT_MUTED),
        (_truncate(activity, 400), 20, False, TEXT_DARK),
    ]
    if duration:
        lines += [("", 10, False, TEXT_DARK), (f"⏱  {duration}", 14, False, TEXT_MUTED)]
    _body_textbox(slide, lines)


def _slide_explain(prs, lesson_data: dict):
    """One flashcard per concept in the Explain phase."""
    concepts = lesson_data.get("explain") or []
    if not concepts:
        slide = _blank_slide(prs)
        _header_band(slide, "explain")
        _body_textbox(slide, [("No concepts provided.", 20, False, TEXT_MUTED)])
        return

    for idx, concept in enumerate(concepts, start=1):
        slide = _blank_slide(prs)
        _header_band(slide, "explain")

        # Concept counter (top-right corner inside band)
        ctr = slide.shapes.add_textbox(SLIDE_W - Inches(1.5), Inches(0.3), Inches(1.2), Inches(0.5))
        p   = ctr.text_frame.paragraphs[0]
        p.text           = f"{idx} / {len(concepts)}"
        p.font.size      = Pt(13)
        p.font.color.rgb = RGBColor(200, 200, 220)
        p.alignment      = PP_ALIGN.RIGHT

        name     = concept.get("name", f"Concept {idx}")
        teaching = concept.get("teaching") or {}
        method   = teaching.get("method", "")
        examples = teaching.get("examples") or []
        duration = teaching.get("duration", "")

        tb = slide.shapes.add_textbox(MARGIN, BODY_TOP, BODY_W, BODY_H)
        tf = tb.text_frame
        tf.word_wrap = True

        # Concept name
        p0 = tf.paragraphs[0]
        p0.text           = name.upper()
        p0.font.bold      = True
        p0.font.size      = Pt(24)
        p0.font.color.rgb = PHASE_META["explain"]["color"]
        p0.space_after    = Pt(10)

        # Teaching method
        if method:
            pm = tf.add_paragraph()
            pm.text           = _truncate(method, 220)
            pm.font.size      = Pt(17)
            pm.font.color.rgb = TEXT_DARK
            pm.space_after    = Pt(10)

        # Examples
        if examples:
            ph = tf.add_paragraph()
            ph.text           = "EXAMPLES"
            ph.font.size      = Pt(11)
            ph.font.bold      = True
            ph.font.color.rgb = TEXT_MUTED
            ph.space_after    = Pt(4)
            for ex in examples[:3]:
                pe = tf.add_paragraph()
                pe.text           = f"•  {ex}"
                pe.font.size      = Pt(15)
                pe.font.color.rgb = TEXT_DARK
                pe.space_after    = Pt(4)

        # Duration
        if duration:
            pd = tf.add_paragraph()
            pd.text           = f"\n⏱  {duration}"
            pd.font.size      = Pt(13)
            pd.font.color.rgb = TEXT_MUTED


def _slide_elaborate(prs, lesson_data: dict):
    elaborate = lesson_data.get("elaborate") or {}
    slide     = _blank_slide(prs)
    _header_band(slide, "elaborate")

    we_do    = elaborate.get("we_do",  "Guided practice activity.")
    you_do   = elaborate.get("you_do", "Independent practice activity.")
    duration = elaborate.get("duration", "")

    # Split the body into two columns
    col_w = (SLIDE_W - MARGIN * 3) / 2
    gap   = MARGIN

    for col_idx, (label, content) in enumerate([("WE DO", we_do), ("YOU DO", you_do)]):
        left = MARGIN + col_idx * (col_w + gap)

        # Coloured label box
        hdr = slide.shapes.add_textbox(left, BODY_TOP, col_w, Inches(0.45))
        ph  = hdr.text_frame.paragraphs[0]
        ph.text           = label
        ph.font.bold      = True
        ph.font.size      = Pt(15)
        ph.font.color.rgb = PHASE_META["elaborate"]["color"]

        # Rule line under label
        rule = slide.shapes.add_shape(6, left, BODY_TOP + Inches(0.45), col_w, Inches(0.03))
        rule.fill.solid()
        rule.fill.fore_color.rgb = PHASE_META["elaborate"]["color"]
        rule.line.fill.background()

        # Content
        body = slide.shapes.add_textbox(left, BODY_TOP + Inches(0.55), col_w, BODY_H - Inches(0.6))
        tf   = body.text_frame
        tf.word_wrap = True
        pc = tf.paragraphs[0]
        pc.text           = _truncate(content, 300)
        pc.font.size      = Pt(17)
        pc.font.color.rgb = TEXT_DARK

    if duration:
        td = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.6), BODY_W, Inches(0.4))
        pd = td.text_frame.paragraphs[0]
        pd.text           = f"⏱  {duration}"
        pd.font.size      = Pt(13)
        pd.font.color.rgb = TEXT_MUTED


def _slide_evaluate(prs, lesson_data: dict):
    evaluate  = lesson_data.get("evaluate") or {}
    questions = evaluate.get("questions") or []
    slide     = _blank_slide(prs)
    _header_band(slide, "evaluate")

    tb = slide.shapes.add_textbox(MARGIN, BODY_TOP, BODY_W, BODY_H)
    tf = tb.text_frame
    tf.word_wrap = True

    if not questions:
        p = tf.paragraphs[0]
        p.text           = "No assessment questions provided."
        p.font.size      = Pt(18)
        p.font.color.rgb = TEXT_MUTED
        return

    first = True
    for q in questions[:6]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first            = False
        p.text           = f"•  {q}"
        p.font.size      = Pt(18)
        p.font.color.rgb = TEXT_DARK
        p.space_after    = Pt(10)


def _slide_closure(prs, lesson_data: dict):
    closure = lesson_data.get("closure")
    if not closure:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "closure")

    summary    = closure.get("summary", "") if isinstance(closure, dict) else str(closure)
    next_topic = closure.get("next_topic", "") if isinstance(closure, dict) else ""

    lines = []
    if summary:
        lines += [("KEY TAKEAWAYS", 12, True, TEXT_MUTED), (_truncate(summary, 350), 18, False, TEXT_DARK)]
    if next_topic:
        lines += [("", 8, False, TEXT_DARK), ("COMING UP NEXT", 12, True, TEXT_MUTED), (next_topic, 18, False, TEXT_DARK)]
    if lines:
        _body_textbox(slide, lines)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

class PPTXService:
    def generate_lesson_pptx(self, lesson_data: dict, output_path: str) -> str:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        _slide_cover(prs, lesson_data)
        _slide_engage(prs, lesson_data)
        _slide_explore(prs, lesson_data)
        _slide_explain(prs, lesson_data)
        _slide_elaborate(prs, lesson_data)
        _slide_evaluate(prs, lesson_data)
        _slide_closure(prs, lesson_data)

        prs.save(output_path)
        return output_path


pptx_service = PPTXService()
